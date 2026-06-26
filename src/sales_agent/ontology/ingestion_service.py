from __future__ import annotations

import json
from pathlib import Path
from typing import Protocol

from sqlalchemy.ext.asyncio import AsyncSession

from sales_agent.llm.base import EmbeddingModel
from sales_agent.models.base import generate_id, utcnow
from sales_agent.models.ingestion import IngestionJob
from sales_agent.ontology.canonicalizer import canonical_key, normalize_aliases
from sales_agent.ontology.conflict import classify_fact_risk, merge_status_for_risk
from sales_agent.ontology.schemas import EntityCandidate, FactCandidate, OntologyIngestionStats


def _docx_text(path: Path) -> str:
    """Extract paragraph text from a .docx (python-docx). Shared by the .docx
    branch and the .doc → .docx conversion path."""
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _doc_to_text(src: Path) -> str:
    """Convert a legacy .doc (OLE2) to .docx via LibreOffice headless, then reuse
    the python-docx paragraph extraction. Chinese/Unicode fidelity is far better
    than antiword/catdoc, which matters because the text feeds the LLM extractor.

    Each call gets its own temp dir + UserInstallation profile so concurrent ingest
    jobs never contend on a shared LibreOffice profile lock."""
    import shutil
    import subprocess
    import tempfile
    if not shutil.which("soffice"):
        raise RuntimeError("未安装 LibreOffice (soffice)，无法解析 .doc")
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp)
        profile_uri = (out / "profile").as_uri()  # absolute file:// URL，soffice 要求
        subprocess.run(
            [
                "soffice", "--headless", "--nologo", "--nofirststartwizard",
                "--norestore", f"-env:UserInstallation={profile_uri}",
                "--convert-to", "docx", "--outdir", str(out), str(src),
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        converted = out / (src.stem + ".docx")
        if not converted.exists():
            raise RuntimeError(f"LibreOffice 转换失败：{src.name}")
        return _docx_text(converted)


def _read_content(path: Path) -> str:
    """Read file content into plain text for the LLM pipeline. Binary office
    formats are parsed by lightweight libs (python-docx / pymupdf / python-pptx /
    openpyxl); legacy .doc is first converted to .docx via LibreOffice headless;
    .md/.txt are read as UTF-8.

    Image files (jpg/png/webp/bmp/gif) are sent to a vision model for AI
    interpretation when ``ontology.vision_enabled`` is True.

    Note: docling was the original choice but its torch-heavy dependency tree
    (2.8GB+) would not install in the target environment (resolver thrash). These
    lightweight libs cover text extraction for the LLM pipeline at a fraction of
    the footprint."""
    ext = path.suffix.lower()
    if ext == ".docx":
        try:
            return _docx_text(path)
        except Exception:
            raise RuntimeError(f"docx 解析失败：{path.name}")
    if ext == ".doc":
        try:
            return _doc_to_text(path)
        except Exception:
            raise RuntimeError(f"doc 解析失败：{path.name}")
    if ext == ".pdf":
        try:
            import fitz  # pymupdf
            parts: list[str] = []
            with fitz.open(str(path)) as doc:
                for page in doc:
                    page_text = page.get_text()
                    if not page_text.strip():
                        # 扫描件/图片页：提取页面为图片，交给视觉模型
                        page_text = _pdf_page_to_vision(page, path.name)
                    parts.append(page_text)
            return "\n".join(parts)
        except Exception:
            raise RuntimeError(f"pdf 解析失败：{path.name}")
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except Exception:
            raise RuntimeError(f"pptx 解析失败：{path.name}")
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filename=str(path), data_only=True, read_only=True)
            texts: list[str] = []
            for ws in wb.worksheets:
                texts.append(f"## {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                    if cells:
                        texts.append("\t".join(cells))
            return "\n".join(texts)
        except Exception:
            raise RuntimeError(f"xlsx 解析失败：{path.name}")
    # ── Image files (vision) ──
    if ext in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}:
        try:
            return _image_to_text_via_vision(path)
        except Exception:
            raise RuntimeError(f"图片解析失败：{path.name}")
    return path.read_text(encoding="utf-8")


def _pdf_page_to_vision(page, filename: str) -> str:
    """Extract text from a scanned/image-based PDF page using vision model.

    Renders the page to a PNG pixmap, encodes as base64, and sends to the
    vision model. Returns the extracted text, or empty string if vision is
    disabled or unavailable.
    """
    from sales_agent.core.config import get_settings
    settings = get_settings()
    if not settings.ontology.vision_enabled:
        return f"[扫描页，视觉解读未开启] {filename}"

    try:
        import base64
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        img_b64 = base64.b64encode(img_bytes).decode("ascii")
        data_url = f"data:image/png;base64,{img_b64}"

        from sales_agent.ontology.img_parser import IMAGE_INTERPRET_PROMPT

        # Use the page-level vision extraction via a lightweight sync wrapper.
        # For production, this should be called from an async context; the
        # synchronous fallback uses httpx directly.
        import httpx
        from sales_agent.core.config import get_settings as _gs
        s = _gs()
        api_key = _get_vision_api_key()
        if not api_key:
            return f"[视觉解读无 API Key] {filename}"

        resp = httpx.post(
            f"{s.model.base_url}/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": s.ontology.vision_model,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": IMAGE_INTERPRET_PROMPT},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }],
                "temperature": 0.1,
                "max_tokens": 1024,
            },
            timeout=60,
        )
        if resp.status_code == 200:
            body = resp.json()
            return body["choices"][0]["message"]["content"].strip()
        return f"[视觉解读 HTTP {resp.status_code}] {filename}"
    except Exception:
        return f"[视觉解读失败] {filename}"


def _image_to_text_via_vision(path: Path) -> str:
    """Extract text from an image file using a vision model.

    Config-driven: requires ``ontology.vision_enabled=True`` and a valid
    API key. Falls back to a placeholder if vision is disabled.
    """
    from sales_agent.core.config import get_settings
    settings = get_settings()
    if not settings.ontology.vision_enabled:
        return f"[图片视觉解读未开启: {path.name}]"

    import httpx
    from sales_agent.ontology.img_parser import (
        IMAGE_INTERPRET_PROMPT,
        get_image_mime_type,
    )

    api_key = _get_vision_api_key()
    if not api_key:
        return f"[视觉解读无 API Key: {path.name}]"

    import base64
    with open(path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode("ascii")

    mime_type = get_image_mime_type(path)
    data_url = f"data:{mime_type};base64,{img_b64}"

    resp = httpx.post(
        f"{settings.model.base_url}/chat/completions",
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        json={
            "model": settings.ontology.vision_model,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": IMAGE_INTERPRET_PROMPT},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }],
            "temperature": 0.1,
            "max_tokens": 2048,
        },
        timeout=60,
    )
    if resp.status_code == 200:
        body = resp.json()
        content = body["choices"][0]["message"]["content"].strip()
        return f"## 图片内容提取: {path.name}\n\n{content}"
    raise RuntimeError(f"视觉解读 HTTP {resp.status_code}: {resp.text[:200]}")


def _get_vision_api_key() -> str:
    """Resolve the API key for vision model calls."""
    import os
    from sales_agent.core.config import get_settings
    settings = get_settings()
    # Check dedicated embedding API key first, then model API key
    env_name = settings.model.embedding_api_key_env or settings.model.api_key_env
    return os.getenv(env_name, "")
    async def extract_entities(self, content: str) -> list[EntityCandidate]: ...
    async def extract_facts(self, content: str, entities: list[EntityCandidate]) -> list[FactCandidate]: ...


class RepositoryProtocol(Protocol):
    async def upsert_entity(self, params: dict) -> tuple[str, bool]: ...
    async def create_fact(self, params: dict) -> str: ...


class OntologyIngestionService:
    def __init__(
        self,
        db: AsyncSession,
        repository: RepositoryProtocol,
        embedding_model: EmbeddingModel,
        extractor: ExtractorProtocol,
    ) -> None:
        self.db = db
        self.repository = repository
        self.embedding_model = embedding_model
        self.extractor = extractor

    async def ingest_paths(
        self,
        *,
        tenant_id: str,
        agent_id: str | None,
        paths: list[Path],
        progress_callback=None,
    ) -> tuple[IngestionJob, OntologyIngestionStats]:
        job = IngestionJob(
            tenant_id=tenant_id,
            agent_id=agent_id,
            engine="ontology_neo4j",
            status="running",
            stage="uploaded",
            documents_seen=len(paths),
        )
        self.db.add(job)
        await self.db.flush()

        stats = OntologyIngestionStats(documents_seen=len(paths))
        for path in paths:
            try:
                await self._ingest_one(job, stats, tenant_id, agent_id, path,
                                       progress_callback=progress_callback)
                stats.documents_ingested += 1
            except Exception as exc:  # noqa: BLE001
                stats.errors.append({"file": str(path), "error": str(exc)})

        job.documents_ingested = stats.documents_ingested
        job.entities_created = stats.entities_created
        job.entities_merged = stats.entities_merged
        job.facts_created = stats.facts_created
        job.facts_active = stats.facts_active
        job.facts_pending_review = stats.facts_pending_review
        job.facts_rejected = stats.facts_rejected
        job.conflicts_created = stats.conflicts_created
        job.warnings_json = json.dumps(stats.warnings, ensure_ascii=False)
        job.errors_json = json.dumps(stats.errors, ensure_ascii=False)
        job.metadata_json = json.dumps(stats.to_metadata(), ensure_ascii=False)
        job.status = "completed" if not stats.errors else "completed_with_errors"
        job.stage = "completed" if not stats.errors else "completed_with_warnings"
        await self.db.flush()
        return job, stats

    async def _ingest_one(
        self,
        job: IngestionJob,
        stats: OntologyIngestionStats,
        tenant_id: str,
        agent_id: str | None,
        path: Path,
        progress_callback=None,
    ) -> None:
        job.stage = "parsed"
        if progress_callback:
            await progress_callback(job.stage, {
                "entities_created": stats.entities_created,
                "entities_merged": stats.entities_merged,
                "facts_created": stats.facts_created,
                "facts_active": stats.facts_active,
                "facts_pending_review": stats.facts_pending_review,
                "conflicts_created": stats.conflicts_created,
            })
        content = _read_content(path)
        source_document_id = generate_id()
        now = utcnow()

        job.stage = "extracting_entities"
        if progress_callback:
            await progress_callback(job.stage, {
                "entities_created": stats.entities_created,
                "entities_merged": stats.entities_merged,
                "facts_created": stats.facts_created,
                "facts_active": stats.facts_active,
                "facts_pending_review": stats.facts_pending_review,
                "conflicts_created": stats.conflicts_created,
            })
        entities = await self.extractor.extract_entities(content)
        embeddings = await self.embedding_model.embed([
            f"{e.name} {json.dumps(e.properties, ensure_ascii=False)}" for e in entities
        ]) if entities else []

        name_to_entity: dict[str, EntityCandidate] = {}
        for entity, embedding in zip(entities, embeddings):
            aliases = normalize_aliases(entity.aliases)
            entity_key = canonical_key(entity.name)
            entity_id, created = await self.repository.upsert_entity({
                "id": generate_id(),
                "tenant_id": tenant_id,
                "agent_id": agent_id,
                "type": entity.type,
                "name": entity.name,
                "canonical_key": entity_key,
                "aliases": aliases,
                "aliases_text": " ".join(aliases),
                "properties": json.dumps(entity.properties, ensure_ascii=False),
                "embedding": embedding,
                "status": "active",
                "now": now,
            })
            name_to_entity[entity.name] = entity
            if created:
                stats.entities_created += 1
            else:
                stats.entities_merged += 1

        job.stage = "extracting_facts"
        if progress_callback:
            await progress_callback(job.stage, {
                "entities_created": stats.entities_created,
                "entities_merged": stats.entities_merged,
                "facts_created": stats.facts_created,
                "facts_active": stats.facts_active,
                "facts_pending_review": stats.facts_pending_review,
                "conflicts_created": stats.conflicts_created,
            })
        facts = await self.extractor.extract_facts(content, entities)
        job.stage = "writing_neo4j"
        if progress_callback:
            await progress_callback(job.stage, {
                "entities_created": stats.entities_created,
                "entities_merged": stats.entities_merged,
                "facts_created": stats.facts_created,
                "facts_active": stats.facts_active,
                "facts_pending_review": stats.facts_pending_review,
                "conflicts_created": stats.conflicts_created,
            })
        for fact in facts:
            subject = name_to_entity.get(fact.subject_name)
            if subject is None:
                stats.warnings.append(f"Fact skipped; subject not found: {fact.subject_name}")
                continue
            object_entity = name_to_entity.get(fact.object_name or "")
            risk = classify_fact_risk(fact)
            status = merge_status_for_risk(risk)
            if status == "active":
                stats.facts_active += 1
            else:
                stats.facts_pending_review += 1
                if risk == "high":
                    stats.conflicts_created += 1
            evidence = fact.evidence[0] if fact.evidence else None
            await self.repository.create_fact({
                "fact_id": generate_id(),
                "tenant_id": tenant_id,
                "agent_id": agent_id,
                "subject_type": subject.type,
                "subject_key": canonical_key(subject.name),
                "object_type": object_entity.type if object_entity else "",
                "object_key": canonical_key(object_entity.name) if object_entity else "",
                "predicate": fact.predicate,
                "fact_type": fact.fact_type,
                "value": fact.value,
                "properties": json.dumps(fact.properties, ensure_ascii=False),
                "confidence": fact.confidence,
                "status": status,
                "risk_level": risk,
                "version": 1,
                "version_date": None,
                "source_document_id": source_document_id,
                "source_title": path.name,
                "source_file_id": None,
                "source_path": str(path),
                "content_hash": "",
                "evidence_id": generate_id(),
                "excerpt": evidence.excerpt if evidence else content[:300],
                "locator": evidence.locator if evidence else path.name,
                "evidence_confidence": evidence.confidence if evidence else fact.confidence,
                "extraction_method": evidence.extraction_method if evidence else "llm",
                "now": now,
            })
            stats.facts_created += 1
